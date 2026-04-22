"""
Coda Knowledge Engine V7.0 — OmniHarvester Plugin.
全格式知识收割器，支持多模态文档 (PDF/Word/Excel/PPT/HTML) 向 KnowledgeNode 的转化。
"""

from __future__ import annotations
import logging
from pathlib import Path
from typing import Any

from ...base_plugin import WikiPlugin, WikiPluginContext, WikiHook
from ...akp_types import KnowledgeNode, NodeType, EpistemicTag

logger = logging.getLogger("Coda.wiki.omni_harvester")


class OmniHarvester(WikiPlugin):
    """
    全格式文档收割器。将二进制或异构格式文档转换为图谱标准的 Markdown Node。
    """
    name = "omni_harvester"

    def __init__(self) -> None:
        self._ctx: WikiPluginContext | None = None
        self.supported_extensions = {".pdf", ".docx", ".xlsx", ".pptx", ".html"}

    async def initialize(self, ctx: WikiPluginContext) -> None:
        self._ctx = ctx
        logger.info("⚙️ OmniHarvester initialized — ready for multi-modal ingestion.")

    async def on_hook(self, hook: WikiHook, payload: Any = None) -> Any:
        # 预留给外部通过 hook 触发解析
        if hook == WikiHook.PRE_COMPILE:
            # 可以在此处执行批量自动收割 (遍历 raw/ 并自动转换为 .md 伴随文件)
            pass
        return None

    def harvest_file(self, filepath: Path) -> list[KnowledgeNode]:
        """将支持的文档格式解析为一个或多个 KnowledgeNode。"""
        if not filepath.exists():
            return []

        ext = filepath.suffix.lower()
        if ext == ".pdf":
            return self._parse_pdf(filepath)
        elif ext == ".docx":
            return self._parse_docx(filepath)
        elif ext == ".pptx":
            return self._parse_pptx(filepath)
        elif ext == ".xlsx":
            return self._parse_xlsx(filepath)
        elif ext in [".html", ".htm"]:
            return self._parse_html(filepath)
        
        return []

    async def extract_knowledge(self, nodes: list[KnowledgeNode]) -> list[KnowledgeNode]:
        """
        [PHASE 2: NER] 对初步解析出的节点内容进行实体与事实提取。
        利用 LLM 将非结构化文本沉淀为高质量的图谱实体和事实。
        """
        if not self._ctx or not self._ctx.llm or not nodes:
            return nodes

        all_nodes = list(nodes)
        for original_node in nodes:
            # 仅解析超过 100 字符且标记为 SOURCE 的节点
            if len(original_node.body) < 100 or original_node.node_type != NodeType.SOURCE:
                continue

            try:
                extraction = await self._run_ner_extraction(original_node.body)
                
                # 1. 提取实体 (Entities)
                for ent in extraction.get("entities", []):
                    ent_node = KnowledgeNode(
                        id=f"ent_{original_node.id}_{ent.get('name', '')[:8]}",
                        title=ent.get("name", "Unknown Entity"),
                        body=ent.get("description", ""),
                        node_type=NodeType.CONCEPT, # 实体标记为概念
                        project_id=original_node.project_id,
                        layer=original_node.layer
                    )
                    # 建立双向链接
                    original_node.links.append(ent_node.id)
                    ent_node.links.append(original_node.id)
                    all_nodes.append(ent_node)

                # 2. 提取事实 (Facts/Claims)
                for fact in extraction.get("facts", []):
                    fact_node = KnowledgeNode(
                        id=f"fact_{original_node.id}_{str(hash(fact))[:8]}",
                        title="Extracted Fact",
                        body=fact,
                        node_type=NodeType.SYNTHESIS, # 事实标记为综合
                        epistemic_tag=EpistemicTag.FACT,
                        project_id=original_node.project_id,
                        layer=original_node.layer
                    )
                    original_node.links.append(fact_node.id)
                    all_nodes.append(fact_node)

            except Exception as e:
                logger.warning(f"NER Extraction failed for node {original_node.id}: {e}")

        return all_nodes

    async def _run_ner_extraction(self, text: str) -> dict[str, Any]:
        """调用 LLM 进行结构化信息抽取。"""
        import json
        prompt = f"""
        你是一个知识发现引擎。请从以下文本中提取关键实体（人、组织、技术、术语）和核心事实/结论。
        输出 JSON 格式。

        文本:
        {text[:2000]} # 截断以节省 token

        要求:
        {{
          "entities": [{{ "name": "...", "description": "..." }}],
          "facts": ["fact 1", "fact 2"]
        }}
        """
        if not self._ctx or not self._ctx.llm:
            return {"entities": [], "facts": []}
        res = await self._ctx.llm.call([{"role": "user", "content": prompt}])
        try:
            # 简单清理和解析 JSON
            import re
            match = re.search(r"\{.*\}", res.text, re.DOTALL)
            if match:
                return json.loads(match.group(0))
        except:
            pass
        return {"entities": [], "facts": []}

    def _create_base_node(self, filepath: Path, title: str, body: str) -> KnowledgeNode:
        import hashlib
        ext = filepath.suffix.lower().lstrip(".")
        node = KnowledgeNode(
            id=f"{filepath.stem}_{hashlib.md5(title.encode()).hexdigest()[:6]}",
            title=title,
            body=body,
            node_type=NodeType.SOURCE,
        )
        # V7.0 联邦字段
        node.source_format = ext
        if self._ctx:
            node.project_id = self._ctx.project_id
            node.layer = self._ctx.layer
            node.readonly = self._ctx.layer <= 1
        return node

    def _parse_pdf(self, filepath: Path) -> list[KnowledgeNode]:
        """使用 PyMuPDF 解析 PDF。"""
        try:
            import fitz  # PyMuPDF
            nodes = []
            doc = fitz.open(str(filepath))
            full_text = []
            for page in doc:
                full_text.append(page.get_text("text"))
            
            body = "\n\n".join(full_text).strip()
            if body:
                nodes.append(self._create_base_node(filepath, filepath.stem, body))
            return nodes
        except BaseException as e:
            logger.error(f"PDF Parse Error ({filepath}): {e}")
            return []

    def _parse_docx(self, filepath: Path) -> list[KnowledgeNode]:
        """使用 python-docx 解析 Word。"""
        try:
            import docx
            doc = docx.Document(str(filepath))
            body = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if body:
                return [self._create_base_node(filepath, filepath.stem, body)]
            return []
        except BaseException as e:
            logger.error(f"DOCX Parse Error ({filepath}): {e}")
            return []

    def _parse_pptx(self, filepath: Path) -> list[KnowledgeNode]:
        """使用 python-pptx 解析 PPT。"""
        try:
            import pptx
            prs = pptx.Presentation(str(filepath))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                text = []
                for shape in slide.shapes:
                    shape_text = getattr(shape, "text", "")
                    if shape_text and shape_text.strip():
                        text.append(shape_text.strip())
                if text:
                    slides_text.append(f"## Slide {i+1}\n" + "\n".join(text))
            
            body = "\n\n".join(slides_text).strip()
            if body:
                return [self._create_base_node(filepath, filepath.stem, body)]
            return []
        except BaseException as e:
            logger.error(f"PPTX Parse Error ({filepath}): {e}")
            return []

    def _parse_xlsx(self, filepath: Path) -> list[KnowledgeNode]:
        """使用 openpyxl 解析 Excel 单元格文本。"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filename=str(filepath), data_only=True)
            sheets_text = []
            for sheet in wb.worksheets:
                rows_text = []
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        rows_text.append(" | ".join(row_data))
                if rows_text:
                    sheets_text.append(f"## Sheet: {sheet.title}\n" + "\n".join(rows_text))
            
            body = "\n\n".join(sheets_text).strip()
            if body:
                return [self._create_base_node(filepath, filepath.stem, body)]
            return []
        except BaseException as e:
            logger.error(f"XLSX Parse Error ({filepath}): {e}")
            return []

    def _parse_html(self, filepath: Path) -> list[KnowledgeNode]:
        """使用 BeautifulSoup4 解析 HTML。"""
        try:
            from bs4 import BeautifulSoup
            content = filepath.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(content, "html.parser")
            title = str(soup.title.string) if soup.title and soup.title.string else filepath.stem
            body = soup.get_text(separator="\n", strip=True)
            if body:
                return [self._create_base_node(filepath, title, body)]
            return []
        except BaseException as e:
            logger.error(f"HTML Parse Error ({filepath}): {e}")
            return []
